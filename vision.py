import os
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import google.generativeai as genai
from PIL import Image
from langchain_community.vectorstores import FAISS
from docx import Document
from pptx import Presentation
# Load API key and configure Google Generative AI
load_dotenv()
genai.configure(api_key="AIzaSyD8zZWCQdyHkQPQOdcQWAOiu_la9-LQ3ZY")
generation_config = {
    "top_p": 0.95,
    "top_k": 40,  # Adjusted to a valid value
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}
# Function to get text from Word, PowerPoint, and PDF files
def get_word_text(word_docs):
    text = ""
    for doc in word_docs:
        document = Document(doc)
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
    return text

def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, index_name):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local(index_name)

def get_conversational_chain():
    prompt_template = """
Based on the context provided, answer the following question in detail. If the answer is not present in the context, state that clearly.
Context:\n {context}\n
Question:\n {question}\n
Answer:
"""

    model = ChatGoogleGenerativeAI(model="gemini-1.5-pro 002", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def handle_pdf_query(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    if not docs:
        return "No relevant documents found. Please ask a different question."
    
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def handle_word_query(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index_word", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    if not docs:
        return "No relevant documents found. Please ask a different question."

    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def handle_ppt_query(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index_ppt", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    if not docs:
        return "No relevant documents found. Please ask a different question."

    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def get_gemini_response(user_input_text, image=None):
    model = genai.GenerativeModel('gemini-1.5-flash')
    chat_session = model.start_chat(history=[])

    if user_input_text:
        if image:
            response = chat_session.send_message([user_input_text, image])
        else:
            response = chat_session.send_message(user_input_text)
    elif image:
        response = chat_session.send_message(image)
    else:
        response = "No input provided."
    return response.text

def main():
    st.set_page_config(page_title="Royce AI", layout="wide")  # Set layout to wide
    st.header("Royce AI 🤖")

    # Initialize session state to hold chat responses
    if 'responses' not in st.session_state:
        st.session_state.responses = []
    if 'user_input_bottom' not in st.session_state:
        st.session_state.user_input_bottom = ""

    

    # Chat history container
    chat_container = st.container()
    with chat_container:
        st.markdown("<div class='chat-container'>", unsafe_allow_html=True)
        for response in st.session_state.responses:
            st.markdown(response)
        st.markdown("</div>", unsafe_allow_html=True)

    # User input box (changed to text_area for multi-line input)
    user_input_bottom = st.text_area("Type your question here...", value=st.session_state.user_input_bottom, placeholder="Ask a question...", key="user_input_bottom", height=100)

    # Expander for uploading documents
    with st.expander("📎 Upload Documents", expanded=False):  # Change 'expanded' to True to keep it open by default
        pdf_docs = st.file_uploader("Upload your PDF Files", type=["pdf"], accept_multiple_files=True, key="pdf_uploader")
        word_docs = st.file_uploader("Upload your Word Files", type=["docx"], accept_multiple_files=True, key="word_uploader")
        ppt_docs = st.file_uploader("Upload your PowerPoint Files", type=["pptx"], accept_multiple_files=True, key="ppt_uploader")
        uploaded_image = st.file_uploader("Upload an image (Optional)", type=["jpeg", "jpg", "png"], key="image_uploader")

    # Function to handle button click and reset input
    def ask_royce():
        if user_input_bottom.strip():  # Check if user input is provided
            # Initialize response variable
            response_text = ""

            # Check for uploaded documents
            if pdf_docs:
                with st.spinner("Searching PDFs..."):
                    raw_text = get_pdf_text(pdf_docs)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks, "faiss_index")  # Specify index name for PDFs
                    st.success("PDFs processed and stored for search.")
                    pdf_response = handle_pdf_query(user_input_bottom)
                    response_text += f"**PDF Response:** {pdf_response}\n"
            if word_docs:
                with st.spinner("Searching Word files..."):
                    raw_text = get_word_text(word_docs)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks, "faiss_index_word")  # Specify index name for Word files
                    st.success("Word files processed and stored for search.")
                    word_response = handle_word_query(user_input_bottom)
                    response_text += f"**Word Response:** {word_response}\n"
            if ppt_docs:
                with st.spinner("Searching PowerPoint files..."):
                    raw_text = get_ppt_text(ppt_docs)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks, "faiss_index_ppt")  # Specify index name for PowerPoint files
                    st.success("PowerPoint files processed and stored for search.")
                    ppt_response = handle_ppt_query(user_input_bottom)
                    response_text += f"**PPT Response:** {ppt_response}\n"
             
            if uploaded_image:
                st.image(uploaded_image, caption="Uploaded Image", use_column_width=True)
                chat_response = get_gemini_response(user_input_bottom, Image.open(uploaded_image))
                response_text += f"**Image Response:** {chat_response}\n"

            # If no responses were generated from documents, handle as a normal chat
            if not response_text:
                chat_response = get_gemini_response(user_input_bottom)
                response_text += f"**Royce's Response:** \n{chat_response}\n"

            # Store responses in session state
            st.session_state.responses.append(f"**Q:** {user_input_bottom}\n\n{response_text}")

            # Clear the input in session state (but before the next render)
            st.session_state.user_input_bottom = ""  # Reset the input
        else:
            st.warning("Please type a question or upload files.")

    # Ask Royce button
    if st.button("Ask Royce", on_click=ask_royce):
        pass  # The function is already handled in the on_click

if __name__ == "__main__":
    main()
